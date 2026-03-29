#!/usr/bin/env python3
"""Generate a sales presentation PPT for Lead Engine Next."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
ACCENT = RGBColor(0x00, 0x96, 0xFF)         # Bright blue
ACCENT2 = RGBColor(0x00, 0xD4, 0xAA)        # Teal green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTITLE_GRAY = RGBColor(0x88, 0x88, 0x88)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_transparency(shape, alpha_pct):
    """Set transparency on a shape. alpha_pct: 0=opaque, 100=fully transparent."""
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    spPr = shape._element.find('.//a:solidFill', nsmap)
    if spPr is not None:
        srgbClr = spPr.find('a:srgbClr', nsmap)
        if srgbClr is not None:
            alpha_elem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        set_shape_transparency(shape, alpha)
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8), icon=""):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {item}" if icon else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
    return txBox


def add_rounded_card(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_icon_number(slide, left, top, number, label, color=ACCENT):
    """Add a big number with label below."""
    card = add_rounded_card(slide, left, top, Inches(2.5), Inches(1.8), RGBColor(0x22, 0x22, 0x3E))
    add_text(slide, left, top + Inches(0.2), Inches(2.5), Inches(0.9), number,
             font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(1.0), Inches(2.5), Inches(0.6), label,
             font_size=14, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 1 - Title
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, PRIMARY)

# Decorative accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

# Title
add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         "Lead Engine", font_size=56, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
         "AI 驱动的全链路智能营销引擎", font_size=32, color=ACCENT)

# Subtitle
add_text(slide, Inches(1.5), Inches(4.0), Inches(8), Inches(0.6),
         "从广告投放 → 线索收集 → 客户运营 → 效果优化，一站式自动化解决方案",
         font_size=18, color=SUBTITLE_GRAY)

# Bottom tagline
add_text(slide, Inches(1.5), Inches(5.8), Inches(5), Inches(0.5),
         "RevoPanda  |  产品介绍", font_size=14, color=SUBTITLE_GRAY)

# Right side decorative element - abstract circles
for i, (x, y, r, c) in enumerate([
    (10.5, 1.0, 2.0, ACCENT),
    (11.0, 3.5, 1.5, ACCENT2),
    (9.5, 4.5, 1.0, ACCENT),
]):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()
    set_shape_transparency(shape, 70 + i * 10)

# =========================================================
# SLIDE 2 - Pain Points
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ORANGE)

add_text(slide, Inches(1.0), Inches(0.5), Inches(5), Inches(0.6),
         "行业痛点", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "传统外贸营销面临的核心挑战", font_size=18, color=SUBTITLE_GRAY)

pain_points = [
    ("01", "广告投放效率低", "手动创建广告素材、设置受众定向、调整预算分配\n耗时耗力，且高度依赖投手个人经验"),
    ("02", "线索流失严重", "客户通过 WhatsApp 咨询后无法及时响应\n大量潜在订单在等待中流失"),
    ("03", "客户跟进断层", "线索分散在多个平台和表格中\n缺乏统一管理和智能跟进机制"),
    ("04", "优化缺乏数据", "广告效果与实际成交之间缺乏闭环追踪\n无法量化 ROI，难以指导优化方向"),
]

for i, (num, title, desc) in enumerate(pain_points):
    row = i // 2
    col = i % 2
    x = Inches(1.0 + col * 6.0)
    y = Inches(2.0 + row * 2.6)

    card = add_rounded_card(slide, x, y, Inches(5.5), Inches(2.2), RGBColor(0x22, 0x22, 0x3E))

    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.8), Inches(0.6),
             num, font_size=28, color=ORANGE, bold=True)
    add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.5),
             title, font_size=20, color=WHITE, bold=True)
    add_text(slide, x + Inches(1.2), y + Inches(0.85), Inches(4), Inches(1.2),
             desc, font_size=13, color=SUBTITLE_GRAY)

# =========================================================
# SLIDE 3 - Solution Overview (Pipeline)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "解决方案总览", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "AI 贯穿全链路，实现营销自动化闭环", font_size=18, color=SUBTITLE_GRAY)

# Pipeline stages
stages = [
    ("市场调研", "竞品分析\n趋势洞察\n受众研究", ACCENT),
    ("方案策划", "预算分配\n受众定向\n平台策略", RGBColor(0x7C, 0x3A, 0xED)),
    ("素材生成", "AI 文案\nAI 图片\n多语言适配", ACCENT2),
    ("广告投放", "一键发布\nMeta Ads\n自动配置", ORANGE),
    ("线索收集", "WhatsApp\n智能对话\n自动提取", RGBColor(0x25, 0xD3, 0x66)),
    ("客户运营", "智能跟进\n质量评估\n自动分配", RGBColor(0xF5, 0xA6, 0x23)),
]

stage_width = Inches(1.7)
gap = Inches(0.15)
total_width = len(stages) * (stage_width + gap) - gap
start_x = (W - total_width) // 2

for i, (title, desc, color) in enumerate(stages):
    x = start_x + i * (stage_width + gap)
    y = Inches(2.2)

    # Card
    card = add_rounded_card(slide, x, y, stage_width, Inches(3.0), RGBColor(0x22, 0x22, 0x3E))

    # Top color bar
    add_shape(slide, x + Inches(0.1), y + Inches(0.1), stage_width - Inches(0.2), Inches(0.08), color)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), y + Inches(0.4), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)

    # Title
    add_text(slide, x, y + Inches(1.15), stage_width, Inches(0.4),
             title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text(slide, x + Inches(0.15), y + Inches(1.6), stage_width - Inches(0.3), Inches(1.2),
             desc, font_size=12, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between stages
    if i < len(stages) - 1:
        arrow_x = x + stage_width
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x - Inches(0.05), Inches(3.5), Inches(0.25), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SUBTITLE_GRAY
        arrow.line.fill.background()

# Bottom highlight
add_rounded_card(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.2), RGBColor(0x22, 0x22, 0x3E))
add_text(slide, Inches(2.0), Inches(5.75), Inches(9), Inches(0.4),
         "核心优势：全流程 AI 自动化  |  对话式操作零门槛  |  数据闭环驱动优化",
         font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(2.0), Inches(6.2), Inches(9), Inches(0.4),
         "无需专业投手经验，通过与 AI 对话即可完成从市场调研到广告上线的全部流程",
         font_size=14, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================
# SLIDE 4 - Campaign Studio (AI Ad Placement)
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "智能广告投放", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "Campaign Studio — 对话式广告创建引擎", font_size=18, color=SUBTITLE_GRAY)

# Left side - Features
features_data = [
    ("对话式广告创建", "通过自然语言描述产品和目标市场\nAI 自动生成完整广告方案"),
    ("AI 市场调研", "自动分析 Meta 广告库竞品素材\n结合 Google Trends 洞察市场趋势"),
    ("智能策略规划", "AI 生成多平台预算分配方案\n精准受众定向与出价策略"),
    ("AI 素材生成", "自动生成多语言广告文案\nAI 图片生成（支持参考图）"),
    ("一键投放上线", "审核确认后一键发布至 Meta\n自动配置广告系列/组/创意"),
]

for i, (title, desc) in enumerate(features_data):
    y = Inches(1.8 + i * 1.05)
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.12), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()

    add_text(slide, Inches(1.6), y, Inches(4.5), Inches(0.35),
             title, font_size=17, color=WHITE, bold=True)
    add_text(slide, Inches(1.6), y + Inches(0.35), Inches(4.5), Inches(0.6),
             desc, font_size=12, color=SUBTITLE_GRAY)

# Right side - Workflow visualization
card = add_rounded_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0), RGBColor(0x22, 0x22, 0x3E))

workflow_steps = [
    ("STEP 1", "需求收集", "描述产品、目标市场、预算", ACCENT),
    ("STEP 2", "市场调研", "竞品分析 + 趋势洞察", RGBColor(0x7C, 0x3A, 0xED)),
    ("STEP 3", "策略生成", "预算分配 + 受众定向", ACCENT2),
    ("STEP 4", "素材创作", "AI 文案 + AI 图片", ORANGE),
    ("STEP 5", "审核发布", "确认 → 一键上线 Meta", RGBColor(0x25, 0xD3, 0x66)),
]

for i, (step, title, desc, color) in enumerate(workflow_steps):
    y = Inches(2.1 + i * 0.9)
    # Step badge
    badge = add_rounded_card(slide, Inches(7.4), y, Inches(1.0), Inches(0.35), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text(slide, Inches(8.6), y - Inches(0.02), Inches(2), Inches(0.35),
             title, font_size=15, color=WHITE, bold=True)
    add_text(slide, Inches(8.6), y + Inches(0.3), Inches(3.5), Inches(0.35),
             desc, font_size=12, color=SUBTITLE_GRAY)

    # Connector line
    if i < len(workflow_steps) - 1:
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.85), y + Inches(0.4), Inches(0.03), Inches(0.45))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = SUBTITLE_GRAY
        line_shape.line.fill.background()

# =========================================================
# SLIDE 5 - Lead Collection
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, RGBColor(0x25, 0xD3, 0x66))

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "智能线索收集", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "WhatsApp AI 客服 — 7×24 小时自动接待与线索提取", font_size=18, color=SUBTITLE_GRAY)

# Left - WhatsApp Bot capabilities
add_text(slide, Inches(1.0), Inches(1.9), Inches(5), Inches(0.5),
         "AI 智能客服能力", font_size=22, color=RGBColor(0x25, 0xD3, 0x66), bold=True)

capabilities = [
    "7×24 全天候自动应答，客户咨询即时响应",
    "支持文字 / 语音 / 图片 / 文档多种消息类型",
    "语音消息自动转文字（Whisper AI 转录）",
    "多产品线智能路由（整车 / 配件 / 农机等）",
    "基于产品知识库的专业回复",
    "复杂需求自动升级至人工销售（飞书通知）",
]

add_bullet_text(slide, Inches(1.2), Inches(2.4), Inches(5.5), Inches(3.5),
                capabilities, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC), icon="▸")

# Right - Lead extraction
add_text(slide, Inches(7.0), Inches(1.9), Inches(5), Inches(0.5),
         "AI 线索自动提取", font_size=22, color=RGBColor(0x25, 0xD3, 0x66), bold=True)

card = add_rounded_card(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.3), RGBColor(0x22, 0x22, 0x3E))

extract_items = [
    ("结构化提取", "从对话中自动识别车型、数量、颜色\n目的国、贸易条款、装运港等关键信息"),
    ("质量分级", "4 级质量评估体系：BAD → GOOD → QUALIFY → PROOF\n自动追问缺失信息，提升线索质量"),
    ("智能去重", "基于车型+目的国自动去重\n防止重复线索干扰销售判断"),
    ("自动归档", "线索自动录入 CRM 数据库\n关联对话记录、广告来源、联系人信息"),
]

for i, (title, desc) in enumerate(extract_items):
    y = Inches(2.7 + i * 1.0)
    add_text(slide, Inches(7.4), y, Inches(4.8), Inches(0.3),
             title, font_size=15, color=WHITE, bold=True)
    add_text(slide, Inches(7.4), y + Inches(0.3), Inches(4.8), Inches(0.65),
             desc, font_size=12, color=SUBTITLE_GRAY)

# =========================================================
# SLIDE 6 - Customer Operations & Optimization
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, RGBColor(0xF5, 0xA6, 0x23))

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "客户运营与效果优化", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "数据闭环驱动持续增长", font_size=18, color=SUBTITLE_GRAY)

# Three columns
col_data = [
    ("客户管理", RGBColor(0xF5, 0xA6, 0x23), [
        "联系人自动建档",
        "多会话历史管理",
        "线索生命周期追踪",
        "人工接管与自动恢复",
        "飞书实时通知销售",
    ]),
    ("数据分析", ACCENT, [
        "对话量趋势分析",
        "线索质量分布统计",
        "询盘国家/车型维度",
        "广告归因追踪",
        "响应时长监控",
    ]),
    ("效果优化", ACCENT2, [
        "广告→线索全链路归因",
        "各广告询盘质量对比",
        "高转化素材识别",
        "预算智能再分配",
        "A/B 测试素材优化",
    ]),
]

for i, (title, color, items) in enumerate(col_data):
    x = Inches(1.0 + i * 4.0)
    y = Inches(1.9)

    card = add_rounded_card(slide, x, y, Inches(3.6), Inches(4.8), RGBColor(0x22, 0x22, 0x3E))

    # Color header bar
    add_shape(slide, x, y, Inches(3.6), Inches(0.6), color)
    add_text(slide, x, y + Inches(0.1), Inches(3.6), Inches(0.4),
             title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        iy = y + Inches(0.9 + j * 0.72)
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), iy + Inches(0.05), Inches(0.2), Inches(0.2))
        check.fill.solid()
        check.fill.fore_color.rgb = color
        check.line.fill.background()
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(9)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        add_text(slide, x + Inches(0.7), iy, Inches(2.7), Inches(0.35),
                 item, font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC))

# =========================================================
# SLIDE 7 - Tech Architecture
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, RGBColor(0x7C, 0x3A, 0xED))

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "技术架构", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "企业级 AI 营销基础设施", font_size=18, color=SUBTITLE_GRAY)

# Architecture layers
layers = [
    ("用户界面层", "Campaign Studio 仪表盘  |  实时 SSE 流式更新  |  多标签工作区", ACCENT, Inches(1.8)),
    ("AI 引擎层", "Claude Sonnet 4.6 策略推理  |  多模型路由  |  AIGC 素材生成  |  Whisper 语音转录", RGBColor(0x7C, 0x3A, 0xED), Inches(2.8)),
    ("业务服务层", "广告编排引擎  |  WhatsApp 消息队列  |  智能客服路由  |  线索提取与评估", ACCENT2, Inches(3.8)),
    ("数据与集成层", "Supabase 数据库  |  Meta Ads API  |  WhatsApp Business API  |  飞书通知", ORANGE, Inches(4.8)),
]

for title, desc, color, y in layers:
    # Layer card
    card = add_rounded_card(slide, Inches(1.0), y, Inches(11.3), Inches(0.85), RGBColor(0x22, 0x22, 0x3E))

    # Left color indicator
    add_shape(slide, Inches(1.0), y, Inches(0.12), Inches(0.85), color)

    add_text(slide, Inches(1.4), y + Inches(0.08), Inches(2.5), Inches(0.35),
             title, font_size=16, color=color, bold=True)
    add_text(slide, Inches(1.4), y + Inches(0.42), Inches(10.5), Inches(0.35),
             desc, font_size=13, color=SUBTITLE_GRAY)

# Key metrics at bottom
add_text(slide, Inches(1.0), Inches(6.0), Inches(3), Inches(0.4),
         "关键技术指标", font_size=18, color=WHITE, bold=True)

metrics = [
    ("24+", "微服务模块"),
    ("5 级", "编排流水线"),
    ("3 层", "LLM 容灾"),
    ("< 2s", "消息响应"),
]

for i, (num, label) in enumerate(metrics):
    x = Inches(1.0 + i * 3.0)
    add_text(slide, x, Inches(6.4), Inches(2.5), Inches(0.5),
             num, font_size=32, color=ACCENT, bold=True)
    add_text(slide, x, Inches(6.9), Inches(2.5), Inches(0.3),
             label, font_size=13, color=SUBTITLE_GRAY)

# =========================================================
# SLIDE 8 - Competitive Advantages
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT2)

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "核心竞争优势", font_size=36, color=WHITE, bold=True)

advantages = [
    ("全链路闭环", "业内唯一覆盖「调研→策划→素材→投放→线索→运营」全链路的 AI 方案\n不是单点工具，而是完整的营销自动化引擎",
     ACCENT, "01"),
    ("对话式操作", "无需学习复杂后台，通过自然语言对话即可完成所有操作\n大幅降低使用门槛，普通业务人员即可上手",
     RGBColor(0x7C, 0x3A, 0xED), "02"),
    ("AI 原生架构", "基于 Claude 4.6 大模型深度集成，不是简单的 API 调用\n具备市场理解、策略推理、创意生成的完整 AI 能力",
     ACCENT2, "03"),
    ("数据驱动优化", "广告投放与线索质量全链路归因\n用真实成交数据反哺广告优化，而非仅依赖平台指标",
     ORANGE, "04"),
]

for i, (title, desc, color, num) in enumerate(advantages):
    row = i // 2
    col = i % 2
    x = Inches(1.0 + col * 6.0)
    y = Inches(1.5 + row * 2.8)

    card = add_rounded_card(slide, x, y, Inches(5.5), Inches(2.4), RGBColor(0x22, 0x22, 0x3E))
    add_shape(slide, x, y, Inches(0.12), Inches(2.4), color)

    add_text(slide, x + Inches(0.4), y + Inches(0.2), Inches(0.8), Inches(0.5),
             num, font_size=28, color=color, bold=True)
    add_text(slide, x + Inches(1.2), y + Inches(0.25), Inches(4), Inches(0.45),
             title, font_size=20, color=WHITE, bold=True)
    add_text(slide, x + Inches(1.2), y + Inches(0.85), Inches(4), Inches(1.3),
             desc, font_size=13, color=SUBTITLE_GRAY)

# =========================================================
# SLIDE 9 - Use Cases / Target Customers
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

add_text(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
         "适用客户", font_size=36, color=WHITE, bold=True)
add_text(slide, Inches(1.0), Inches(1.1), Inches(8), Inches(0.5),
         "为出海企业量身打造", font_size=18, color=SUBTITLE_GRAY)

customers = [
    ("外贸整车出口", "面向非洲/中东/南美市场的整车出口商\n高客单价、长决策周期、线索质量至关重要",
     ACCENT, "🚗"),
    ("汽车配件出口", "SKU 多样、询盘频繁\n需要快速响应和精准产品匹配",
     ACCENT2, "🔧"),
    ("农业机械出口", "目标市场分散、产品专业性强\n需要多语言、多市场同步投放",
     ORANGE, "🚜"),
    ("跨境电商品牌", "需要规模化广告投放\n追求高效线索转化和 ROI 优化",
     RGBColor(0x7C, 0x3A, 0xED), "🌍"),
]

for i, (title, desc, color, icon) in enumerate(customers):
    x = Inches(1.0 + i * 3.0)
    y = Inches(2.0)

    card = add_rounded_card(slide, x, y, Inches(2.7), Inches(4.5), RGBColor(0x22, 0x22, 0x3E))

    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y + Inches(0.3), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    set_shape_transparency(circle, 75)

    add_text(slide, x, y + Inches(0.4), Inches(2.7), Inches(0.8),
             icon, font_size=36, alignment=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(1.6), Inches(2.7), Inches(0.5),
             title, font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.3), Inches(2.0),
             desc, font_size=12, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================
# SLIDE 10 - CTA / Contact
# =========================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)

# Decorative circles
for x, y, r, c in [(10.0, 0.5, 2.5, ACCENT), (11.5, 4.0, 2.0, ACCENT2)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(r), Inches(r))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()
    set_shape_transparency(shape, 90)

add_text(slide, Inches(1.5), Inches(2.0), Inches(8), Inches(1.0),
         "让 AI 驱动您的出海增长", font_size=48, color=WHITE, bold=True)

add_text(slide, Inches(1.5), Inches(3.3), Inches(8), Inches(0.6),
         "从第一条广告到第一笔订单，Lead Engine 为您全程护航",
         font_size=20, color=SUBTITLE_GRAY)

# CTA button
btn = add_rounded_card(slide, Inches(1.5), Inches(4.5), Inches(3.5), Inches(0.8), ACCENT)
tf = btn.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = "预约产品演示 →"
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(6)

add_text(slide, Inches(1.5), Inches(5.8), Inches(8), Inches(0.5),
         "RevoPanda  |  AI-Powered Marketing Engine",
         font_size=16, color=SUBTITLE_GRAY)

# Save
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           'Lead_Engine_产品介绍.pptx')
prs.save(output_path)
print(f"PPT saved to: {output_path}")
